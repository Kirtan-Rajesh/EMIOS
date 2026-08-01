"""Document Processing Service: extract -> chunk -> embed -> store, for RAG
(UML Phase 3 / Phase 5's data side). Handles unstructured documents (PDF, DOCX,
plain text) attached to an assessment. Structured CSV/JSON uploads keep using
the existing legacy ingestion path into the graph
(app/api/endpoints.py's POST /api/upload) - this service is specifically for
making free-text company documents searchable via RAG, a separate concern.

Blocking work (PDF/DOCX parsing, embedding API calls) is offloaded to Starlette's
threadpool via run_in_threadpool so it doesn't stall the asyncio event loop -
the async DB write at the end still uses the normal async SQLAlchemy session.
"""

from __future__ import annotations

import asyncio
import io
import logging
import threading
from concurrent.futures import ThreadPoolExecutor
from typing import List

from starlette.concurrency import run_in_threadpool

from app.core.embeddings import get_embedding_provider
from app.entities.document_chunk import DocumentChunk
from app.repositories.document_chunk_repository import DocumentChunkRepository
from app.services.vector_search_service import get_vector_store

logger = logging.getLogger("emios")

# Fixed-size sliding-window chunking with overlap - not a dedicated splitter
# library. Simple, dependency-free, and easy to reason about; fine for the
# document sizes this project deals with.
CHUNK_SIZE_CHARS = 1500
CHUNK_OVERLAP_CHARS = 200

# At most this many embedding-API calls are ever in flight at once, process-
# wide - not just per document. Bounded rather than unbounded so a large
# document (or several documents being processed concurrently, see
# UploadService.upload_zip_archive_stream) don't collectively fire dozens of
# simultaneous requests at the provider and trip Bedrock's per-account TPS
# throttling. A plain threading.Semaphore (not per-call) is what makes this
# a *global* cap regardless of how many callers are embedding at once.
_MAX_CONCURRENT_EMBEDS = 8
_embed_semaphore = threading.BoundedSemaphore(_MAX_CONCURRENT_EMBEDS)

_PDF_CONTENT_TYPES = {"application/pdf"}
_DOCX_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
_PPTX_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
_XLSX_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


class DocumentProcessingService:
    def __init__(self, chunk_repo: DocumentChunkRepository):
        self.chunk_repo = chunk_repo

    def extract_text(self, file_bytes: bytes, content_type: str, filename: str) -> str:
        """Extracts plain text from PDF, DOCX, PPTX, XLSX, or anything else treated
        as plain text (txt/csv/json/markdown). HLDs in particular are
        frequently authored as slide decks, hence PPTX support here; asset/app
        inventories are frequently authored as spreadsheets, hence XLSX."""
        lower_name = filename.lower()
        if content_type in _PDF_CONTENT_TYPES or lower_name.endswith(".pdf"):
            return self._extract_pdf_text(file_bytes)
        if content_type in _DOCX_CONTENT_TYPES or lower_name.endswith(".docx"):
            return self._extract_docx_text(file_bytes)
        if content_type in _PPTX_CONTENT_TYPES or lower_name.endswith(".pptx"):
            return self._extract_pptx_text(file_bytes)
        if content_type in _XLSX_CONTENT_TYPES or lower_name.endswith(".xlsx"):
            return self._extract_xlsx_text(file_bytes)
        return file_bytes.decode("utf-8", errors="replace")

    @staticmethod
    def _extract_pdf_text(file_bytes: bytes) -> str:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    @staticmethod
    def _extract_docx_text(file_bytes: bytes) -> str:
        import docx

        document = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    @staticmethod
    def _extract_pptx_text(file_bytes: bytes) -> str:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(file_bytes))
        lines: List[str] = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_lines = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if slide_lines:
                lines.append(f"[Slide {slide_num}]")
                lines.extend(slide_lines)
        return "\n".join(lines)

    @staticmethod
    def _extract_xlsx_text(file_bytes: bytes) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines: List[str] = []
        for sheet in workbook.worksheets:
            sheet_lines = [
                "\t".join("" if cell is None else str(cell) for cell in row)
                for row in sheet.iter_rows(values_only=True)
                if any(cell is not None and str(cell).strip() for cell in row)
            ]
            if sheet_lines:
                lines.append(f"[Sheet: {sheet.title}]")
                lines.extend(sheet_lines)
        return "\n".join(lines)

    def chunk_text(self, text: str) -> List[str]:
        """Splits text into overlapping fixed-size chunks. Empty/whitespace-only
        text produces zero chunks (not an error - just nothing to index)."""
        text = text.strip()
        if not text:
            return []
        chunks = []
        start = 0
        while start < len(text):
            end = start + CHUNK_SIZE_CHARS
            chunks.append(text[start:end])
            if end >= len(text):
                break
            start = end - CHUNK_OVERLAP_CHARS
        return chunks

    def embed_chunks(self, chunks: List[str]) -> List[List[float]]:
        """Embeds every chunk concurrently instead of one at a time - each
        embed() call is a network round trip to the provider (Bedrock/
        OpenAI), so it's pure I/O wait and parallelizes well; a document with
        10 chunks previously meant 10 sequential round trips. Actual
        concurrency is capped process-wide by _embed_semaphore (not just
        within this call) - documents can be processed concurrently too (see
        UploadService.upload_zip_archive_stream), and the semaphore is what
        keeps the *total* number of simultaneous Bedrock calls bounded no
        matter how many documents are in flight at once.
        ThreadPoolExecutor.map() preserves input order in its output, so
        embeddings[i] still matches chunks[i]; if any chunk fails to embed,
        the exception propagates and the whole document fails, same as the
        previous sequential behavior (no partial-chunk documents)."""
        provider = get_embedding_provider()

        def embed_one(text: str) -> List[float]:
            with _embed_semaphore:
                return provider.embed(text)

        with ThreadPoolExecutor(max_workers=min(16, len(chunks))) as pool:
            return list(pool.map(embed_one, chunks))

    async def process(
        self,
        assessment_id: str,
        upload_id: str,
        file_bytes: bytes,
        content_type: str,
        filename: str,
        db_lock: asyncio.Lock | None = None,
    ) -> int:
        """Runs the full extract -> chunk -> embed -> store pipeline for one
        upload. Returns the number of chunks stored.

        `db_lock`, if given, is held only around the final bulk_create() -
        the extract/chunk/embed/upsert steps above it run fully unlocked.
        Passed by UploadService.upload_zip_archive_stream(), where several
        documents are processed concurrently via `self.chunk_repo`s each
        bound to their own short-lived session sharing one underlying
        connection in tests (SQLite + StaticPool - see tests/conftest.py);
        overlapping writes on that shared connection corrupt each other's
        transaction state, so those specific writes are serialized. Callers
        processing one document at a time (the default, None) don't need
        this and pay no locking overhead."""
        text = await run_in_threadpool(self.extract_text, file_bytes, content_type, filename)
        chunks = self.chunk_text(text)
        if not chunks:
            logger.info(f"No extractable text in upload '{upload_id}' ({filename}) - 0 chunks stored.")
            return 0

        embeddings = await run_in_threadpool(self.embed_chunks, chunks)
        # get_vector_store() constructs a QdrantVectorStore, whose __init__ does a
        # synchronous get_collections()/create_collection() round trip on first use -
        # must go through the threadpool too, not just the upsert call below.
        vector_store = await run_in_threadpool(get_vector_store)
        point_ids = await run_in_threadpool(
            vector_store.upsert_chunks, assessment_id, upload_id, chunks, embeddings
        )

        chunk_rows = [
            DocumentChunk(
                upload_id=upload_id,
                assessment_id=assessment_id,
                chunk_index=i,
                text=chunks[i],
                qdrant_point_id=point_ids[i],
            )
            for i in range(len(chunks))
        ]
        if db_lock is not None:
            async with db_lock:
                await self.chunk_repo.bulk_create(chunk_rows)
        else:
            await self.chunk_repo.bulk_create(chunk_rows)
        return len(chunks)
